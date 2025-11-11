from pptx import Presentation

prs = Presentation('Sample-MidTerm-Presentation-TermProject.pptx')

for i, slide in enumerate(prs.slides):
    print(f'\n=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if shape.text.strip():
                print(shape.text)
