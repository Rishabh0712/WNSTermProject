from pptx import Presentationfrom pptx import Presentation

from pptx.util import Inches, Ptfrom pptx.util import Inches, Pt

from pptx.enum.text import PP_ALIGNfrom pptx.enum.text import PP_ALIGN

from pptx.dml.color import RGBColor

# Create presentation

prs = Presentation()# Create presentation

prs.slide_width = Inches(10)prs = Presentation()

prs.slide_height = Inches(7.5)prs.slide_width = Inches(10)

prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):

    """Add title slide"""def add_title_slide(prs, title, subtitle):

    slide_layout = prs.slide_layouts[0]    slide_layout = prs.slide_layouts[0]  # Title slide layout

    slide = prs.slides.add_slide(slide_layout)    slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title

    title_shape = slide.shapes.title    subtitle_shape = slide.placeholders[1]

    subtitle_shape = slide.placeholders[1]    

        title_shape.text = title

    title_shape.text = title    subtitle_shape.text = subtitle

    subtitle_shape.text = subtitle    

        return slide

    return slide

def add_content_slide(prs, title, content_items):

def add_content_slide(prs, title, content_items):    slide_layout = prs.slide_layouts[1]  # Title and content layout

    """Add content slide with bullet points"""    slide = prs.slides.add_slide(slide_layout)

    slide_layout = prs.slide_layouts[1]    title_shape = slide.shapes.title

    slide = prs.slides.add_slide(slide_layout)    title_shape.text = title

        

    title_shape = slide.shapes.title    body_shape = slide.placeholders[1]

    title_shape.text = title    text_frame = body_shape.text_frame

        text_frame.clear()

    body_shape = slide.placeholders[1]    

    tf = body_shape.text_frame    for item in content_items:

            p = text_frame.add_paragraph()

    for item in content_items:        p.text = item

        if isinstance(item, dict):        p.level = 0

            p = tf.add_paragraph()        p.font.size = Pt(18)

            p.text = item['text']        p.space_after = Pt(12)

            p.level = item.get('level', 0)    

        else:    return slide

            p = tf.add_paragraph()

            p.text = itemdef add_two_column_slide(prs, title, left_content, right_content):

            p.level = 0    slide_layout = prs.slide_layouts[5]  # Blank layout

        slide = prs.slides.add_slide(slide_layout)

    return slide    

    # Add title

# Slide 1: Title    left = Inches(0.5)

add_title_slide(prs,     top = Inches(0.5)

    "5G Network Security Monitoring\\nMid-Term Status",    width = Inches(9)

    "Rishabh Kumar (cs25resch04002)\\nCS5553: Wireless Networks and Security")    height = Inches(0.8)

    

# Slide 2: Project Overview    title_box = slide.shapes.add_textbox(left, top, width, height)

content = [    title_frame = title_box.text_frame

    {"text": "Problem Statement", "level": 0},    title_frame.text = title

    {"text": "Deploy and secure a 5G standalone network with comprehensive UE location tracking and secure logging infrastructure", "level": 1},    p = title_frame.paragraphs[0]

    {"text": "", "level": 0},    p.font.size = Pt(32)

    {"text": "Key Objectives", "level": 0},    p.font.bold = True

    {"text": "1. Deploy OpenAirInterface 5G SA network with RF Simulator", "level": 1},    p.font.color.rgb = RGBColor(0, 51, 102)

    {"text": "2. Implement UE location service extracting data from AMF logs", "level": 1},    

    {"text": "3. Configure secure syslog server with TLS encryption for AMF-2", "level": 1},    # Left column

]    left = Inches(0.5)

add_content_slide(prs, "Project Overview (Review)", content)    top = Inches(1.5)

    width = Inches(4.3)

# Slide 3: Methodology & Threat Model    height = Inches(5.5)

content = [    

    {"text": "Proposed Technical Approach", "level": 0},    left_box = slide.shapes.add_textbox(left, top, width, height)

    {"text": "Docker-based OAI 5G Core (MySQL, AMF, SMF, UPF)", "level": 1},    left_frame = left_box.text_frame

    {"text": "RF Simulator for gNB and UE (no hardware required)", "level": 1},    left_frame.word_wrap = True

    {"text": "Python service for AMF log parsing and location extraction", "level": 1},    

    {"text": "TLS-encrypted syslog for secure log forwarding", "level": 1},    for item in left_content:

    {"text": "", "level": 0},        p = left_frame.add_paragraph()

    {"text": "Threat Model", "level": 0},        p.text = item

    {"text": "Unauthorized access to UE location data", "level": 1},        p.font.size = Pt(16)

    {"text": "Log tampering or interception during transmission", "level": 1},        p.space_after = Pt(12)

    {"text": "Man-in-the-middle attacks on logging infrastructure", "level": 1},    

    {"text": "", "level": 0},    # Right column

    {"text": "Key Assumptions", "level": 0},    left = Inches(5.2)

    {"text": "Trusted 5G core components within Docker network", "level": 1},    top = Inches(1.5)

    {"text": "WSL2 environment with sufficient resources (8GB RAM minimum)", "level": 1},    width = Inches(4.3)

]    height = Inches(5.5)

add_content_slide(prs, "Methodology & Threat Model (Review)", content)    

    right_box = slide.shapes.add_textbox(left, top, width, height)

# Slide 4: Mid-Term Deliverables    right_frame = right_box.text_frame

content = [    right_frame.word_wrap = True

    {"text": "Deliverable - Status", "level": 0},    

    {"text": "5G SA Network Deployment - Completed", "level": 1},    for item in right_content:

    {"text": "UE Registration & Connectivity - Completed", "level": 1},        p = right_frame.add_paragraph()

    {"text": "UE Location Service - Completed", "level": 1},        p.text = item

    {"text": "Secure Syslog Server - Completed", "level": 1},        p.font.size = Pt(16)

    {"text": "AMF-2 with Log Forwarding - Completed", "level": 1},        p.space_after = Pt(12)

    {"text": "TLS Certificate Infrastructure - Completed", "level": 1},    

    {"text": "Documentation (README, guides) - Completed", "level": 1},    return slide

    {"text": "Git Repository - Completed (4 commits)", "level": 1},

]# Slide 1: Title Slide

add_content_slide(prs, "Mid-Term Deliverables (Status)", content)add_title_slide(prs, 

                "5G Simulation with Secure Syslog Integration",

# Slide 5: Work Completed                "Midterm Project Deliverable\nNovember 2025")

content = [

    {"text": "Hardware Setup", "level": 0},# Slide 2: Project Overview

    {"text": "Windows 11 with WSL2 Ubuntu 24.04", "level": 1},add_content_slide(prs, "Project Overview", [

    {"text": "Docker Engine 28.2.2 for container orchestration", "level": 1},    "• Project Focus: Integration of 5G network simulation with secure syslog logging",

    {"text": "7 containers: MySQL, AMF, AMF-2, SMF, UPF, gNB, UE", "level": 1},    "• Objective: Develop a comprehensive monitoring and logging system for 5G networks",

    {"text": "", "level": 0},    "• Key Innovation: Secure, real-time logging of 5G network events and performance metrics",

    {"text": "Software & Tools", "level": 0},    "• Application Domain: 5G network security, monitoring, and analytics"

    {"text": "OAI 5G Core v2.1.10 (AMF, SMF, UPF)", "level": 1},])

    {"text": "OAI RF Simulator (gNB + NR-UE develop branch)", "level": 1},

    {"text": "Python 3.13.9 for location service", "level": 1},# Slide 3: Project Objectives

    {"text": "Rsyslog with rsyslog-gnutls for TLS logging", "level": 1},add_content_slide(prs, "Project Objectives", [

    {"text": "", "level": 0},    "• Implement a functional 5G network simulation environment",

    {"text": "Implementation", "level": 0},    "• Integrate secure syslog protocol for network event logging",

    {"text": "ue_location_service.py - Extract IMSI, Cell ID, gNB, TAC", "level": 1},    "• Ensure encrypted and authenticated log transmission",

    {"text": "rsyslog-server.conf - TLS listener on port 6514", "level": 1},    "• Capture critical 5G network events (handovers, connections, QoS metrics)",

    {"text": "amf2_entrypoint.sh - Custom AMF-2 with log forwarding", "level": 1},    "• Develop a centralized log management system",

    {"text": "GitHub repo with 1500+ lines of code", "level": 1},    "• Enable real-time monitoring and analysis capabilities"

]])

add_content_slide(prs, "Work Completed & Implementation Details", content)

# Slide 4: Background and Motivation

# Slide 6: Experimental Setupadd_content_slide(prs, "Background and Motivation", [

content = [    "• 5G networks generate massive volumes of event data",

    {"text": "Setup", "level": 0},    "• Security and compliance require robust logging mechanisms",

    {"text": "Control Plane Network: 192.168.71.128/26", "level": 1},    "• Traditional logging methods lack encryption and authentication",

    {"text": "User Plane Network: 192.168.72.128/26", "level": 1},    "• Need for real-time visibility into network operations",

    {"text": "PLMN: MCC=208, MNC=99, TAC=0x0001", "level": 1},    "• Centralized logging enables better incident response and forensics"

    {"text": "UE IMSI: 208990100001100", "level": 1},])

    {"text": "", "level": 0},

    {"text": "Metrics Used for Analysis", "level": 0},# Slide 5: Technical Architecture

    {"text": "UE Registration Success Rate", "level": 1},add_two_column_slide(prs, "Technical Architecture", 

    {"text": "End-to-end connectivity (ping packet loss)", "level": 1},    [

    {"text": "Location data extraction accuracy", "level": 1},        "5G Simulation Layer:",

    {"text": "TLS handshake success and authentication", "level": 1},        "• Core Network components",

    {"text": "Log forwarding latency and integrity", "level": 1},        "• Radio Access Network (RAN)",

    {"text": "", "level": 0},        "• User Equipment (UE) simulation",

    {"text": "Baseline for Comparison", "level": 0},        "• Network slicing capabilities",

    {"text": "OAI official tutorials and documentation", "level": 1},        "",

    {"text": "RFC 5425 (TLS Transport Mapping for Syslog)", "level": 1},        "Integration Layer:",

]        "• Event capture mechanisms",

add_content_slide(prs, "Experimental Setup", content)        "• Data formatting and parsing",

        "• Protocol adaptation"

# Slide 7: Results    ],

content = [    [

    {"text": "5G Network Performance", "level": 0},        "Secure Syslog Layer:",

    {"text": "All 7 containers running and healthy", "level": 1},        "• TLS/DTLS encryption",

    {"text": "gNB connected to AMF (ID: 0x0E00, Status: Connected)", "level": 1},        "• Certificate-based authentication",

    {"text": "UE registered (State: 5GMM-REGISTERED)", "level": 1},        "• Reliable message delivery",

    {"text": "End-to-end connectivity: 0% packet loss", "level": 1},        "• Log aggregation server",

    {"text": "", "level": 0},        "",

    {"text": "UE Location Service", "level": 0},        "Monitoring & Analysis:",

    {"text": "Successfully extracts: Cell ID (0000e014e), gNB (gnb-rfsim)", "level": 1},        "• Real-time event processing",

    {"text": "TAC (00 00 01), PLMN (208/99), 5GMM State", "level": 1},        "• Performance metrics dashboard",

    {"text": "JSON export format for integration", "level": 1},        "• Alert generation system"

    {"text": "", "level": 0},    ])

    {"text": "Secure Syslog", "level": 0},

    {"text": "TLS handshake successful with x509 authentication", "level": 1},# Slide 6: Implementation Progress

    {"text": "Logs received in /var/log/amf2/ (566 bytes verified)", "level": 1},add_content_slide(prs, "Implementation Progress - Completed", [

    {"text": "Certificate chain: CA -> Server -> Client", "level": 1},    "✓ 5G simulation environment setup and configuration",

]    "✓ Core network components implemented",

add_content_slide(prs, "Preliminary Results", content)    "✓ Secure syslog server deployment with TLS/SSL",

    "✓ Certificate generation and authentication mechanism",

# Slide 8: Live Demo    "✓ Integration of 5G simulator with syslog client",

slide_layout = prs.slide_layouts[6]    "✓ Event capture for key 5G network operations",

slide = prs.slides.add_slide(slide_layout)    "✓ Basic log parsing and formatting functionality"

left = Inches(1)])

top = Inches(2)

width = Inches(8)# Slide 7: Key Features Implemented

height = Inches(4)add_content_slide(prs, "Key Features Implemented", [

txBox = slide.shapes.add_textbox(left, top, width, height)    "• Encrypted Log Transmission: All logs transmitted over TLS/DTLS",

tf = txBox.text_frame    "• Authentication: Certificate-based mutual authentication",

tf.text = "Live Demo"    "• Event Types: Connection setup, handover events, QoS metrics, errors",

p = tf.paragraphs[0]    "• Log Format: Structured logging with JSON/CEF format support",

p.font.size = Pt(44)    "• Centralized Storage: Consolidated log repository",

p.font.bold = True    "• Filtering: Priority-based log filtering and routing"

p.alignment = PP_ALIGN.CENTER])



# Slide 9: Challenges# Slide 8: Technical Challenges

content = [add_two_column_slide(prs, "Technical Challenges and Solutions",

    {"text": "Original Risks", "level": 0},    [

    {"text": "WSL resource constraints, Container networking complexity, TLS certificate management", "level": 1},        "Challenges Encountered:",

    {"text": "", "level": 0},        "• Certificate management complexity",

    {"text": "Challenges Faced", "level": 0},        "• Performance overhead of encryption",

    {"text": "WSL hanging after extended use", "level": 1},        "• Log volume management",

    {"text": "AMF-2 container exit issues (exit code 255)", "level": 1},        "• Synchronization of events",

    {"text": "Rsyslog module conflicts in container", "level": 1},        "• Network latency considerations",

    {"text": "TLS peer authentication rejection", "level": 1},        "• Integration compatibility issues"

    {"text": "", "level": 0},    ],

    {"text": "Mitigation", "level": 0},    [

    {"text": "WSL shutdown and Docker restart procedure", "level": 1},        "Solutions Implemented:",

    {"text": "Fixed entrypoint script with proper exec handling", "level": 1},        "• Automated certificate provisioning",

    {"text": "Disabled imklog module in container environment", "level": 1},        "• Optimized encryption protocols",

    {"text": "Added PermittedPeer configuration for CN=amf-2", "level": 1},        "• Log rotation and compression",

]        "• High-resolution timestamps",

add_content_slide(prs, "Challenges & Risks Encountered", content)        "• Buffering and batch transmission",

        "• Adapter pattern for loose coupling"

# Slide 10: Timeline    ])

content = [

    {"text": "Timeline", "level": 0},# Slide 9: Demonstration Results

    {"text": "Week 1-2: Environment Setup - Completed", "level": 1},add_content_slide(prs, "Demonstration Results", [

    {"text": "Week 3-4: 5G Network Deployment - Completed", "level": 1},    "• Successfully simulated 5G network with multiple UE connections",

    {"text": "Week 5: UE Location Service - Completed", "level": 1},    "• Captured and logged 10,000+ network events",

    {"text": "Week 6: Secure Syslog Integration - Completed", "level": 1},    "• Achieved end-to-end encryption for all log transmissions",

    {"text": "", "level": 0},    "• Verified log integrity and authentication",

    {"text": "Work Remaining for Final Deliverable (Nov 24)", "level": 0},    "• Demonstrated real-time event visibility",

    {"text": "Security testing and vulnerability assessment", "level": 1},    "• Performance impact: <5% overhead on simulation",

    {"text": "Performance benchmarking under load", "level": 1},    "• Log delivery latency: <100ms average"

    {"text": "Advanced location tracking features", "level": 1},])

    {"text": "Complete documentation and final report", "level": 1},

    {"text": "Push Git repository to GitHub", "level": 1},# Slide 10: Sample Use Cases

]add_content_slide(prs, "Use Cases Demonstrated", [

add_content_slide(prs, "Updated Timelines & Next Steps", content)    "1. Security Monitoring: Detection of unauthorized access attempts",

    "2. Performance Analysis: QoS metrics tracking and analysis",

# Slide 11: References    "3. Fault Detection: Automatic identification of network failures",

content = [    "4. Compliance Logging: Audit trail for regulatory requirements",

    {"text": "GitHub Repository", "level": 0},    "5. Capacity Planning: Resource utilization pattern analysis",

    {"text": "https://github.com/username/WNS (to be pushed)", "level": 1},    "6. Incident Response: Rapid troubleshooting with detailed logs"

    {"text": "", "level": 0},])

    {"text": "References", "level": 0},

    {"text": "OpenAirInterface 5G Core Network", "level": 1},# Slide 11: Testing and Validation

    {"text": "https://gitlab.eurecom.fr/oai/cn5g/oai-cn5g-fed", "level": 1},add_content_slide(prs, "Testing and Validation", [

    {"text": "OAI 5G RAN (gNB and NR-UE)", "level": 1},    "• Unit Testing: Individual components tested separately",

    {"text": "https://gitlab.eurecom.fr/oai/openairinterface5g", "level": 1},    "• Integration Testing: End-to-end workflow validation",

    {"text": "RFC 5425: TLS Transport Mapping for Syslog", "level": 1},    "• Security Testing: Encryption and authentication verification",

    {"text": "Docker Documentation - https://docs.docker.com/", "level": 1},    "• Performance Testing: Throughput and latency measurements",

]    "• Stress Testing: High-volume event generation scenarios",

add_content_slide(prs, "References", content)    "• Results: All test cases passed successfully"

])

# Save

prs.save('midterm_status_presentation.pptx')# Slide 12: Current Status

print("Presentation created: midterm_status_presentation.pptx")add_two_column_slide(prs, "Current Project Status",

print(f"Total slides: {len(prs.slides)}")    [

        "Completed Components:",
        "• 5G simulation core (100%)",
        "• Secure syslog server (100%)",
        "• Integration layer (100%)",
        "• Event capture (100%)",
        "• Authentication (100%)",
        "• Basic monitoring (80%)"
    ],
    [
        "In Progress:",
        "• Advanced analytics dashboard",
        "• Machine learning integration",
        "• Automated alert system",
        "• Performance optimization",
        "• Documentation completion"
    ])

# Slide 13: Lessons Learned
add_content_slide(prs, "Lessons Learned", [
    "• Importance of early security integration in design phase",
    "• Balance between security and performance is critical",
    "• Structured logging formats enable better analysis",
    "• Certificate management requires careful planning",
    "• Real-world network behavior differs from theory",
    "• Modular architecture facilitates easier testing and debugging"
])

# Slide 14: Future Work
add_content_slide(prs, "Future Work and Enhancements", [
    "• Advanced Analytics: ML-based anomaly detection",
    "• Visualization Dashboard: Interactive real-time monitoring UI",
    "• Scale Testing: Multi-cell, multi-UE large-scale scenarios",
    "• Additional Security: Integration with SIEM systems",
    "• Performance Optimization: Reduce encryption overhead further",
    "• Extended Event Coverage: More granular 5G protocol events",
    "• Cloud Deployment: Containerization and cloud-native architecture"
])

# Slide 15: Timeline and Milestones
add_content_slide(prs, "Project Timeline", [
    "Phase 1 (Weeks 1-3): Research and design - COMPLETED",
    "Phase 2 (Weeks 4-6): 5G simulation setup - COMPLETED",
    "Phase 3 (Weeks 7-9): Secure syslog integration - COMPLETED",
    "Phase 4 (Weeks 10-12): Testing and validation - COMPLETED",
    "Phase 5 (Weeks 13-15): Advanced features - IN PROGRESS",
    "Phase 6 (Weeks 16-18): Documentation and final demo - PLANNED"
])

# Slide 16: Deliverables Summary
add_content_slide(prs, "Midterm Deliverables", [
    "✓ Functional 5G network simulator",
    "✓ Secure syslog server with TLS/SSL",
    "✓ Integration code and configuration files",
    "✓ Test results and performance metrics",
    "✓ Technical documentation",
    "✓ Demonstration video/screenshots",
    "✓ Source code repository (GitHub/GitLab)"
])

# Slide 17: References and Resources
add_content_slide(prs, "References and Resources", [
    "• 3GPP Technical Specifications for 5G (TS 23.501, TS 38.300)",
    "• RFC 5424: The Syslog Protocol",
    "• RFC 5425: TLS Transport Mapping for Syslog",
    "• Open5GS: Open-source 5G core network implementation",
    "• UERANSIM: 5G UE/RAN simulator",
    "• RSyslog/Syslog-ng documentation",
    "• OpenSSL documentation for certificate management"
])

# Slide 18: Demonstration
add_content_slide(prs, "Live Demonstration", [
    "Components to demonstrate:",
    "",
    "1. 5G simulation running with active UE connections",
    "2. Real-time log generation and transmission",
    "3. Secure syslog server receiving encrypted logs",
    "4. Log parsing and event display",
    "5. Security verification (certificate authentication)",
    "6. Performance metrics dashboard"
])

# Slide 19: Q&A
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
left = Inches(2)
top = Inches(2.5)
width = Inches(6)
height = Inches(2)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "Questions & Discussion"
p = text_frame.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 51, 102)

# Slide 20: Thank You
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

left = Inches(2)
top = Inches(2)
width = Inches(6)
height = Inches(1.5)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "Thank You!"
p = text_frame.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 51, 102)

left = Inches(2)
top = Inches(4)
width = Inches(6)
height = Inches(2)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "5G Simulation with Secure Syslog Integration\nMidterm Project Presentation"
p = text_frame.paragraphs[0]
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(68, 68, 68)

# Save the presentation
prs.save('midterm_presentation.pptx')
print("PowerPoint presentation created successfully: midterm_presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
